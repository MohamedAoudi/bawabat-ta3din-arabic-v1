"""Validate AMIP PPTX report templates against the placeholder contract."""

from __future__ import annotations

import argparse
import re
import sys
from collections import Counter
from pathlib import Path
from typing import Iterable

try:
    from pptx import Presentation
except ImportError as exc:  # pragma: no cover - exercised only without dependency
    raise SystemExit(
        "python-pptx is required to validate PPTX templates. "
        "Install it with: pip install -r requirements.txt"
    ) from exc


PLACEHOLDER_PATTERN = re.compile(r"\{\{[A-Z][A-Z0-9_]*\}\}")

OPTIONAL_PLACEHOLDERS = frozenset(
    {
        "{{PRICE_NARRATIVE}}",
        "{{CHART_TOP_IMPORT_PARTNERS}}",
        "{{CHART_HS_BREAKDOWN}}",
        "{{TABLE_HS_BREAKDOWN}}",
        "{{KPI_TOP_HS_PRODUCT}}",
    }
)

REQUIRED_PLACEHOLDERS = frozenset(
    {
        "{{REPORT_TITLE}}",
        "{{REPORT_SUBTITLE}}",
        "{{COUNTRY_NAME}}",
        "{{MINERAL_NAME}}",
        "{{YEAR_RANGE}}",
        "{{LANGUAGE}}",
        "{{GENERATED_AT}}",
        "{{DATA_COVERAGE_NOTE}}",
        "{{EXECUTIVE_SUMMARY}}",
        "{{INTRODUCTION_TEXT}}",
        "{{MINERAL_IMPORTANCE_TEXT}}",
        "{{REPORT_OBJECTIVES_TEXT}}",
        "{{KPI_LATEST_PRODUCTION}}",
        "{{KPI_LATEST_PRODUCTION_YEAR}}",
        "{{KPI_PRODUCTION_CHANGE}}",
        "{{KPI_EXPORT_VALUE}}",
        "{{KPI_IMPORT_VALUE}}",
        "{{KPI_TRADE_BALANCE}}",
        "{{KPI_TOP_EXPORT_PARTNER}}",
        "{{KPI_TOP_IMPORT_PARTNER}}",
        "{{CHART_PRODUCTION_TREND}}",
        "{{CHART_EXPORT_IMPORT_TREND}}",
        "{{CHART_TRADE_BALANCE}}",
        "{{CHART_TOP_EXPORT_PARTNERS}}",
        "{{CHART_PRODUCTION_RANKING}}",
        "{{TABLE_KPI_SUMMARY}}",
        "{{TABLE_PRODUCTION_SERIES}}",
        "{{TABLE_TRADE_SERIES}}",
        "{{TABLE_PARTNER_TRADE}}",
        "{{TABLE_STRENGTHS_CHALLENGES}}",
        "{{TABLE_ANNEX_PRODUCTION}}",
        "{{TABLE_ANNEX_TRADE}}",
        "{{PRODUCTION_NARRATIVE}}",
        "{{TRADE_NARRATIVE}}",
        "{{PARTNER_TRADE_NARRATIVE}}",
        "{{HS_BREAKDOWN_NARRATIVE}}",
        "{{COMPARATIVE_ANALYSIS}}",
        "{{RISK_FLAGS}}",
        "{{OPPORTUNITY_FLAGS}}",
        "{{CONCLUSIONS}}",
        "{{RECOMMENDATIONS}}",
        "{{DATA_SOURCES}}",
    }
)

KNOWN_PLACEHOLDERS = REQUIRED_PLACEHOLDERS | OPTIONAL_PLACEHOLDERS


def _iter_shapes(shapes: Iterable[object]) -> Iterable[object]:
    """Yield shapes recursively, including shapes nested inside groups."""
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from _iter_shapes(shape.shapes)


def _extract_placeholders(text: str) -> list[str]:
    return PLACEHOLDER_PATTERN.findall(text or "")


def scan_pptx_placeholders(pptx_path: Path) -> list[str]:
    """Return placeholders found in text boxes and table cells, preserving repeats."""
    presentation = Presentation(str(pptx_path))
    placeholders: list[str] = []

    for slide in presentation.slides:
        for shape in _iter_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                placeholders.extend(_extract_placeholders(shape.text))

            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        placeholders.extend(_extract_placeholders(cell.text))

    return placeholders


def validate_pptx_template(pptx_path: Path) -> dict[str, list[str]]:
    """Validate a PPTX template and return found, missing, and unknown placeholders."""
    found = scan_pptx_placeholders(pptx_path)
    found_unique = sorted(set(found))

    return {
        "found": found_unique,
        "missing_required": sorted(REQUIRED_PLACEHOLDERS - set(found)),
        "missing_optional": sorted(OPTIONAL_PLACEHOLDERS - set(found)),
        "unknown": sorted(set(found) - KNOWN_PLACEHOLDERS),
        "duplicates": sorted(
            placeholder
            for placeholder, count in Counter(found).items()
            if count > 1
        ),
    }


def print_validation_summary(pptx_path: Path, result: dict[str, list[str]]) -> None:
    print("AMIP PPTX template validation")
    print(f"Template: {pptx_path}")
    print(f"Found placeholders: {len(result['found'])}")
    print(f"Required placeholders: {len(REQUIRED_PLACEHOLDERS)}")
    print(f"Optional placeholders: {len(OPTIONAL_PLACEHOLDERS)}")
    print()

    if result["found"]:
        print("Found placeholders:")
        for placeholder in result["found"]:
            print(f"  - {placeholder}")
    else:
        print("Found placeholders: none")

    print()
    if result["missing_required"]:
        print("Missing required placeholders:")
        for placeholder in result["missing_required"]:
            print(f"  - {placeholder}")
    else:
        print("Missing required placeholders: none")

    print()
    if result["missing_optional"]:
        print("Missing optional placeholders:")
        for placeholder in result["missing_optional"]:
            print(f"  - {placeholder}")
    else:
        print("Missing optional placeholders: none")

    print()
    if result["unknown"]:
        print("Unknown placeholders (warnings):")
        for placeholder in result["unknown"]:
            print(f"  - {placeholder}")
    else:
        print("Unknown placeholders: none")

    print()
    if result["duplicates"]:
        print("Duplicate placeholders (warnings):")
        for placeholder in result["duplicates"]:
            print(f"  - {placeholder}")
    else:
        print("Duplicate placeholders: none")

    print()
    if result["missing_required"]:
        print("Validation failed: required placeholders are missing.")
    else:
        print("Validation succeeded: all required placeholders are present.")


def parse_args(argv: list[str] | None = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Validate an AMIP PPTX template against the placeholder contract."
    )
    parser.add_argument(
        "pptx_path",
        type=Path,
        help="Path to the PPTX template to validate.",
    )
    return parser.parse_args(argv)


def main(argv: list[str] | None = None) -> int:
    args = parse_args(argv)
    pptx_path = args.pptx_path

    if not pptx_path.exists():
        print(f"Template not found: {pptx_path}", file=sys.stderr)
        return 2

    if not pptx_path.is_file():
        print(f"Template path is not a file: {pptx_path}", file=sys.stderr)
        return 2

    result = validate_pptx_template(pptx_path)
    print_validation_summary(pptx_path, result)
    return 1 if result["missing_required"] else 0


if __name__ == "__main__":
    raise SystemExit(main())
